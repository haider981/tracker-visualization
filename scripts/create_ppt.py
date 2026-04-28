"""
Create a PowerPoint presentation explaining the Work Tracker Analytics visualization.
Uses Editorial_Maths team as the example for the first view. Includes all screenshots
with arrows and explanatory text.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Paths
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)
ASSETS_DIR = os.path.join(PROJECT_ROOT, 'presentation_assets')
OUTPUT_PATH = os.path.join(PROJECT_ROOT, 'Work_Tracker_Analytics_Presentation.pptx')

# Ordered list of images matching slide flow (Dashboard view first, then Project view)
# Based on user's 17 screenshots - order follows logical explanation flow
IMAGE_FILES = [
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-771bbecb-549a-419c-a7fc-4d8ac3cfbc81.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-5e06e2d8-4ee2-4249-84ff-a6d0d0981da6.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-486e6f15-8453-44d3-a542-0413b2ca7252.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-72b81b4a-9b0a-4cb4-bb55-a700b12d91f5.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-b934b56c-ce7d-4c20-988e-183a6aea99b0.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-0ff0ae1c-3d23-46ac-9ee7-5b0bb007c1b6.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-e6f95fcb-be8b-4c24-9e77-a290ad0fb4e4.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-742123b9-4fb6-4723-b907-d2fba2382592.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-1a1e0950-13c3-4c44-842e-26ec0be313be.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-185d757b-2377-4f37-8a21-aef64dd93966.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-cd949020-90b4-4539-9d78-47d93394c320.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-724cbee5-5422-4c7a-a053-eb8610ff5bb9.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-d9f80778-b94d-41f7-a856-2a456a80f85b.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-ef5e778d-8cbc-4b8e-9f74-4a6e360c5432.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-e193fe18-2f74-4219-a90b-d70bc7f85e13.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-25f7f601-56d4-4097-9309-b6eaae2e73de.png',
    'c__Users_Haider_AppData_Roaming_Cursor_User_workspaceStorage_images_image-22efb5dc-bda0-4468-8425-e3f48593afa4.png',
]

# Slide configurations: title, annotation text, arrows = [(left, top, width, height, label), ...]
SLIDES = [
    {"title": "Dashboard Overview - Key Metrics & Charts", "text": "① FILTERS (top): Department, Team, Employee, Period - use to narrow data\n② KPI CARDS: Total Projects, Active Employees, Total Hours, Total Tasks\n③ TOP 10 PROJECTS: Bar chart - click bar for book elements breakdown\n④ WORK MODE ANALYSIS: Stacked bars show days per employee by work mode (In Office, WFH, etc.)"},
    {"title": "Task Types, Task Status & Team Contributions", "text": "① TASK TYPES: Pie chart - distribution of task categories (COM, CR1, DRF, R1, etc.)\n② TASK STATUS: Bar chart - Completed vs In Progress tasks\n③ TEAM CONTRIBUTIONS: Each team shows Hours (teal) and Units (orange) - compare productivity"},
    {"title": "Project Activity Timeline, Work Mode & Book Elements", "text": "① PROJECT ACTIVITY TIMELINE: Stacked area chart - work hours over time per project\n② WORK MODE DISTRIBUTION: Pie chart - % of time in each mode (In Office, WFH, etc.)\n③ BOOK ELEMENTS: Bar chart - Chapter, Exercise, Full book, etc. (hours/count per element)"},
    {"title": "Filters - Department Dropdown (Example: Editorial)", "text": "① DEPARTMENT: Select All, DTP, Editorial, or Digital Marketing\n② TEAM NAME: Filters to teams in selected department (e.g., Editorial_Maths)\n③ EMPLOYEE: Search to filter by specific person\n④ PERIOD: All, Last 7/30/90 days, Last 6 Months, etc.\n⑤ CLEAR ALL: Reset all filters"},
    {"title": "Filters - Selecting Editorial_Maths Team", "text": "① TEAM NAME dropdown: Select 'Editorial_Maths' to view only that team's data\n② All charts and metrics below update automatically for Editorial_Maths\n③ Employee list shows members: Jaspal Singh, Kirtish Kumar, Mahesh Kumawat, etc."},
    {"title": "Editorial_Maths Team - Filtered Metrics", "text": "With Department=Editorial, Team=Editorial_Maths:\n① TOTAL PROJECTS: 56 (projects this team works on)\n② ACTIVE EMPLOYEES: 11 (team size)\n③ TOTAL TASKS: 4516\n④ TOP 10 PROJECTS: Bar chart shows this team's top projects by hours\n⑤ WORK MODE ANALYSIS: Stacked bars for each Editorial_Maths employee"},
    {"title": "Employee Drill-down: Ashok Kumar Sharma", "text": "Filter by Employee = Ashok Kumar Sharma:\n① TOTAL PROJECTS: 15 | ACTIVE EMPLOYEES: 1 | TOTAL HOURS: 1614\n② TOP 10 PROJECTS: His top projects by hours (CBSE_AppMath_SM highest)\n③ WORK MODE: One stacked bar - In Office (blue) dominant, OT Home, Night, Leave, etc.\n④ Click bar for book elements; hover for details"},
    {"title": "Top 10 Projects - Book Elements Tooltip", "text": "When you CLICK a project bar (e.g., FK_8th_CBSE_Math_SM):\n① TOOLTIP shows: Total hours + Book Elements breakdown\n② Example: Chapter 1278.5 hrs, Exercise 125 hrs, Full book 79 hrs, Diagram 41.5 hrs\n③ Total Tasks: 512 for that project\n④ Enables drill-down from project → element-level analysis"},
    {"title": "Work Mode Analysis - Employee Breakdown (Anshika Srivastava)", "text": "Click/hover on an employee's bar for detailed breakdown:\n① IN OFFICE: 112 days - Top 5 projects with hours\n② OT HOME: 29 days - Top 5 projects\n③ WFH: 14 days - Top 5 projects\n④ LEAVE: 12 days\n⑤ TOTAL: 182 days - combines all modes"},
    {"title": "Project Activity Timeline - Editorial_Maths (Hover Tooltip)", "text": "With Editorial_Maths filter:\n① TIMELINE shows work hours per day, stacked by project\n② HOVER on a date: Tooltip shows Total hrs + per-project breakdown\n③ Example: 2025-08-20 → Total 68 hrs, VK_7th_CBSE_Math 38.5, FK_7th 20, etc.\n④ Legend below: color = project"},
    {"title": "Work Mode & Book Elements - Editorial_Maths", "text": "① WORK MODE PIE: 72.9% In Office, 14.3% other, etc. (Editorial_Maths only)\n② BOOK ELEMENTS BAR: Chapter ~5500, Exercise ~700, etc.\n③ HOVER on Exercise bar: Tooltip shows Task Types (DRF, CMPL-MS, R3, R1, R2) with hours\n④ Drill-down: Element → Task Types"},
    {"title": "Task Types & Status - Editorial_Maths (R1 Tooltip)", "text": "① TASK TYPES PIE: Hover R1 → Tooltip shows Total Hours (1283.25) + Top 5 Book Elements\n   (R1→Chapter 946.75 hrs, R1→Exercise 81 hrs, etc.) + Total Units 4296\n② TASK STATUS: Completed (purple) vs In Progress (pink) bars\n③ Enables: Task Type → Book Element drill-down"},
    {"title": "Team Contributions - Editorial_Maths Tooltip", "text": "Hover on Editorial_Maths bar in Team Contributions:\n① TOTAL: 86405.26 hrs (combined metric)\n② UNITS: 77143 units (orange bar - output measure)\n③ HOURS: 9262.26 hrs (teal bar - time logged)\n④ TOTAL TASKS: 2611\n⑤ Compare Hours vs Units across teams for productivity"},
    {"title": "Project Activity Timeline - Editorial_Maths (Last 6 Months)", "text": "Period filter = Last 6 Months, Team = Editorial_Maths:\n① Stacked area shows daily work hours Aug 2025 - Feb 2026\n② Each color = different project (VK_7th_CBSE_Math, FK_7th, etc.)\n③ Legend scrolls horizontally - many projects\n④ Use to see activity trends and peak periods"},
    {"title": "Project View - Overview", "text": "Different view: Project-wise task breakdown & timeline\n① FILTERS: Segment, Series, Class, Period, Project (multi-select with search)\n② PROJECT-WISE TASKS: Horizontal stacked bars - each row = project, segments = task types (DRF, Coord, CR, R1, etc.)\n③ PROJECT TIMELINE: Line graph - activity over time"},
    {"title": "Project View - Series & Project Filters", "text": "① SERIES dropdown: LCreate, LDC, CodeSpark, ComW11, FKComp, etc. - product lines\n② PROJECT search: Type to find projects, checkboxes for multi-select\n③ SEGMENT, CLASS, PERIOD: Additional filters\n④ Task legend: colors map to task types in stacked bars"},
    {"title": "Project View - Timeline with Daily Tooltip", "text": "Hover on Project Timeline for a specific date:\n① TOOLTIP: Date (e.g., 2025-11-27)\n② Task | Hours breakdown: Analysis 7.5, MEET 5, Training 4, Coord 3, Research 1.5\n③ Total: 21 hrs for that day\n④ Each line color = task type - track activity over time"},
]

def add_arrow(slide, left, top, width=0.4, height=0.2):
    """Add orange arrow shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 165, 0)
    shape.line.color.rgb = RGBColor(255, 165, 0)
    return shape

def add_callout(slide, left, top, width, height, text, font_size=10):
    """Add callout text box with dark bg."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = False
    p.font.color.rgb = RGBColor(255, 255, 255)
    # Add fill
    sp = box._element
    spPr = sp.get_or_add_spPr()
    from pptx.oxml.ns import qn
    solidFill = spPr.get_or_add_solidFill()
    srgbClr = solidFill.get_or_add_srgbClr()
    srgbClr.set(qn('val'), '2D3748')
    return box

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Apply dark background (optional - via slide background)
    # Blank layout 6

    # ========== SLIDE 1: Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "Work Tracker Analytics"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(139, 92, 246)

    sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(12), Inches(0.8))
    p2 = sub.text_frame.paragraphs[0]
    p2.text = "Complete Visualization Guide"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(200, 200, 200)

    desc = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1.8))
    p3 = desc.text_frame.paragraphs[0]
    p3.text = "This presentation explains every feature of the Work Tracker Analytics dashboard.\nUsing Editorial_Maths team as the primary example for clarity."
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(160, 160, 160)

    # ========== SLIDE 2: Introduction ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.7))
    t.text_frame.paragraphs[0].text = "What is Work Tracker Analytics?"
    t.text_frame.paragraphs[0].font.size = Pt(28)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = RGBColor(139, 92, 246)

    bullet = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
    tf = bullet.text_frame
    tf.word_wrap = True
    points = [
        "• Real-time insights into work activity across departments and teams",
        "• Track projects, hours, tasks, and employee work modes",
        "• Two main views: Dashboard View (team/employee focused) and Project View (project focused)",
        "• Filter by: Department, Team Name, Employee, Period",
        "• Auto-refreshes every minute for up-to-date data",
    ]
    for i, pt in enumerate(points):
        p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.space_after = Pt(10)

    # ========== SLIDES 3-19: Screenshots with annotations ==========
    for i, config in enumerate(SLIDES):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.55))
        title_box.text_frame.paragraphs[0].text = config["title"]
        title_box.text_frame.paragraphs[0].font.size = Pt(20)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(139, 92, 246)

        img_path = None
        if i < len(IMAGE_FILES):
            img_path = os.path.join(ASSETS_DIR, IMAGE_FILES[i])
            if not os.path.exists(img_path):
                img_path = None

        if img_path and os.path.exists(img_path):
            # Add image - fit width to 12 inches
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(0.72), width=Inches(12))
            # Annotation box
            ann = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12), Inches(1.0))
            ann.text_frame.word_wrap = True
            ann.text_frame.paragraphs[0].text = config["text"]
            ann.text_frame.paragraphs[0].font.size = Pt(11)
            ann.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
        else:
            # Fallback: text only
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(5))
            tb.text_frame.word_wrap = True
            tb.text_frame.paragraphs[0].text = config["text"]
            tb.text_frame.paragraphs[0].font.size = Pt(16)
            tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

    # ========== Final slide ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    t = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(9), Inches(1))
    t.text_frame.paragraphs[0].text = "Thank You"
    t.text_frame.paragraphs[0].font.size = Pt(52)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = RGBColor(139, 92, 246)
    t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(9), Inches(1))
    sub.text_frame.paragraphs[0].text = "Questions?"
    sub.text_frame.paragraphs[0].font.size = Pt(24)
    sub.text_frame.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)
    sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return prs

def main():
    if not os.path.exists(ASSETS_DIR):
        print(f"Assets folder not found: {ASSETS_DIR}")
        print("Please ensure presentation_assets folder exists with screenshot images.")
    prs = create_presentation()
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    return OUTPUT_PATH

if __name__ == '__main__':
    main()
